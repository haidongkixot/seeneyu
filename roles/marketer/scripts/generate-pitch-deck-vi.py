#!/usr/bin/env python3
"""
seeneyu Vietnamese Pitch Deck Generator
Generates: .shared/outputs/marketer/seeneyu-pitch-deck-vi.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
DARK_BG = RGBColor(0x0D, 0x0D, 0x14)
DARK_BG_HEX = "0D0D14"
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_DARK = RGBColor(0xD9, 0x87, 0x06)
WHITE = RGBColor(0xE8, 0xE8, 0xF0)
GREY = RGBColor(0x8F, 0x8F, 0xAD)
DARK_SURFACE = RGBColor(0x1A, 0x1A, 0x2E)
DARK_CARD = RGBColor(0x16, 0x16, 0x24)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

pptx = Presentation()
pptx.slide_width = SLIDE_W
pptx.slide_height = SLIDE_H

# ── Helpers ──

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_text(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font=FONT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = align
    return txBox

def add_multiline(slide, lines, x, y, w, h, size=16, color=WHITE, line_spacing=1.5, bold=False):
    """Add multiple lines of text in one text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.space_after = Pt(size * (line_spacing - 1))
    return txBox

def add_bullet_list(slide, items, x, y, w, h, size=15, color=WHITE, bullet_color=AMBER):
    """Add a bulleted list with amber bullets."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Use bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25CF  "
        run_bullet.font.size = Pt(size - 2)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = FONT

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(size)
        run_text.font.color.rgb = color
        run_text.font.name = FONT
        p.space_after = Pt(8)
    return txBox

def add_rect(slide, x, y, w, h, fill_color=DARK_SURFACE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_accent_line(slide, x, y, w):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AMBER
    shape.line.fill.background()
    return shape

def slide_title(slide, title, subtitle=None):
    add_text(slide, title, 0.8, 0.4, 10, 0.7, size=32, color=AMBER, bold=True)
    add_accent_line(slide, 0.8, 1.1, 2.5)
    if subtitle:
        add_text(slide, subtitle, 0.8, 1.3, 10, 0.5, size=16, color=GREY, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Trang bia (Cover)
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])  # blank
set_bg(s)

# Decorative amber accent bar at top
add_rect(s, 0, 0, 13.333, 0.08, AMBER)

# Logo / App name
add_text(s, "seeneyu", 1.5, 1.8, 10, 1.2, size=72, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Tagline
add_text(s, "Huan luyen giao tiep bang AI", 1.5, 3.2, 10, 0.7, size=28, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(s, "Xem \u2192 Quan sat \u2192 Bat chuoc \u2192 Phan hoi AI \u2192 Lap lai", 1.5, 4.0, 10, 0.5, size=18, color=GREY, align=PP_ALIGN.CENTER)

# URL
add_text(s, "seeneyu.vercel.app", 1.5, 5.0, 10, 0.4, size=16, color=GREY, align=PP_ALIGN.CENTER)

# Amber box for seed round
add_rect(s, 4.5, 5.8, 4.3, 0.7, AMBER)
add_text(s, "SEED ROUND \u2014 $1.5M", 4.5, 5.85, 4.3, 0.6, size=22, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

# Company
add_text(s, "PeeTeeAI JSC  |  www.peetees.ai", 1.5, 6.7, 10, 0.4, size=13, color=GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: Van de
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Van de")

add_text(s, "55% thong diep truyen tai qua ngon ngu co the", 0.8, 1.8, 11, 0.7, size=28, color=AMBER, bold=True)
add_text(s, "\u2014 Albert Mehrabian, 1967", 0.8, 2.4, 6, 0.4, size=14, color=GREY, italic=True)

bullets = [
    "Ky nang giao tiep phi ngon ngu rat quan trong nhung hau nhu khong duoc day bai ban",
    "Doc sach khong tao phan xa co bap \u2014 can thuc hanh lap di lap lai",
    "Xem lai ban than tren Zoom khong co cau truc, khong co phan hoi cu the",
    "Khong co nen tang nao chuyen ve luyen tap ngon ngu co the",
    "Lam viec tu xa khien giao tiep phi ngon ngu kho doc hon bao gio het",
]
add_bullet_list(s, bullets, 0.8, 3.0, 11, 3.5, size=17)

# Stats box
add_rect(s, 8.5, 5.8, 4, 1.2, DARK_SURFACE, AMBER)
add_text(s, "$370 TY", 8.5, 5.85, 4, 0.5, size=26, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "chi cho dao tao doanh nghiep toan cau", 8.5, 6.35, 4, 0.5, size=13, color=GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: Giai phap
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Giai phap: seeneyu")

add_text(s, "Nen tang huan luyen giao tiep bang AI dau tien", 0.8, 1.6, 11, 0.6, size=22, color=WHITE, bold=True)
add_text(s, "su dung Hollywood lam hinh mau", 0.8, 2.1, 11, 0.5, size=22, color=AMBER, bold=True)

# Flow diagram boxes
steps = [
    ("1. XEM", "Xem clip tu phim\nkinh dien"),
    ("2. QUAN SAT", "Phan tich ky thuat\nqua Observation Guide"),
    ("3. BAT CHUOC", "Quay video ban than\nluyen tap ky nang"),
    ("4. PHAN HOI AI", "GPT-4o Vision\ncham diem chi tiet"),
    ("5. LAP LAI", "Thuc hanh den khi\nthanh thao"),
]

x_start = 0.5
for i, (title, desc) in enumerate(steps):
    x = x_start + i * 2.5
    add_rect(s, x, 3.2, 2.2, 2.2, DARK_SURFACE, AMBER if i == 3 else None)
    add_text(s, title, x, 3.3, 2.2, 0.5, size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.1, 3.9, 2.0, 1.2, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < 4:
        add_text(s, "\u25B6", x + 2.2, 3.9, 0.3, 0.5, size=18, color=AMBER, align=PP_ALIGN.CENTER)

add_text(s, '"Think Duolingo cho ngon ngu co the, duoc GPT-4o Vision ho tro"', 0.8, 5.8, 11, 0.5, size=17, color=GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: Cach hoat dong
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Cach hoat dong")

# Left column - Core Loop
add_rect(s, 0.8, 1.8, 5.5, 5.0, DARK_SURFACE)
add_text(s, "VONG LAP HOC TAP", 0.8, 1.9, 5.5, 0.5, size=18, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

loop_items = [
    "1. Duyet thu vien 100+ bai hoc tu phim kinh dien",
    "2. Xem clip YouTube voi huong dan quan sat chi tiet",
    "3. Luyen tap theo buoc nho (Micro-Practice, 30 giay moi buoc)",
    "4. Quay video qua webcam, AI phan tich khung hinh",
    "5. Nhan diem so + goi y cu the de cai thien",
    "6. Lam lai cho den khi thanh thao \u2192 chuyen bai tiep theo",
]
add_bullet_list(s, loop_items, 1.0, 2.5, 5.0, 4.0, size=14, color=WHITE)

# Right column - Arcade
add_rect(s, 7.0, 1.8, 5.5, 5.0, DARK_SURFACE)
add_text(s, "ARCADE \u2014 THU THACH NHANH", 7.0, 1.9, 5.5, 0.5, size=18, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

arcade_items = [
    "Thu thach bieu cam khuon mat trong 10 giay",
    "Thu thach cu chi co the theo mau tham chieu",
    "GPT-4o Vision so sanh va cham diem 0\u2013100",
    "30+ thu thach chia thanh cac bo de (bundles)",
    "Phu hop de khoi dong hoac luyen nhanh moi ngay",
    "Gamification: diem so, tien trinh, phan thuong",
]
add_bullet_list(s, arcade_items, 7.2, 2.5, 5.0, 4.0, size=14, color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: Tinh nang chinh
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Tinh nang chinh")

features = [
    ("Foundation Lessons", "3 khoa hoc nen tang (Giong noi, Ngon tu, Ngon ngu co the)\nmoi khoa 10+ bai, ly thuyet + video YouTube + quiz"),
    ("Thu vien 100+ Bai hoc", "Clip tu phim kinh dien, phan loai theo 5 ky nang\nObservation Guide + Micro-Practice Steps cho moi clip"),
    ("Arcade Challenges", "Thu thach nhanh 10 giay \u2014 bieu cam & cu chi\nCham diem AI so sanh voi hinh mau tham chieu"),
    ("AI Coaching Loop", "GPT-4o Vision phan tich khung hinh video\nPhan hoi chi tiet: diem so + ke hoach hanh dong cu the"),
    ("Lo trinh ca nhan hoa", "Danh gia 5 ky nang khi onboarding\nHe thong goi y bai hoc theo trinh do"),
    ("Admin CMS", "Quan ly clip, crawl YouTube tu dong, AI scoring\nImport ZIP, quan ly arcade, analytics dashboard"),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.8 + row * 2.7
    add_rect(s, x, y, 3.9, 2.3, DARK_CARD, None)
    add_text(s, title, x + 0.2, y + 0.15, 3.5, 0.45, size=16, color=AMBER, bold=True)
    add_text(s, desc, x + 0.2, y + 0.7, 3.5, 1.4, size=12, color=GREY)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: Cong nghe
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Cong nghe")

add_text(s, "Ket hop AI tien tien nhat de tao trai nghiem huan luyen doc nhat", 0.8, 1.6, 11, 0.5, size=18, color=WHITE)

tech_stack = [
    ("GPT-4o Vision", "Phan tich ngon ngu co the\ntu khung hinh video", "AI Core"),
    ("Whisper STT", "Nhan dien giong noi\ncho bai tap vocal", "Voice"),
    ("OpenAI TTS", "Phat am phan hoi\ncua Coach Ney", "Voice"),
    ("MediaPipe", "ML phia client\nphat hien khuon mat/cu chi", "Client ML"),
    ("Next.js 14", "App Router, TypeScript\nServer Components", "Framework"),
    ("Neon PostgreSQL", "Co so du lieu\nqua Prisma ORM", "Database"),
    ("Vercel Blob", "Luu tru video\nguoi dung tam thoi", "Storage"),
    ("YouTube API", "Thu vien video clip\nkhong tu host", "Content"),
]

for i, (name, desc, category) in enumerate(tech_stack):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.1
    y = 2.4 + row * 2.4
    add_rect(s, x, y, 2.9, 2.0, DARK_SURFACE)
    add_text(s, category, x + 0.15, y + 0.1, 2.6, 0.3, size=10, color=GREY)
    add_text(s, name, x + 0.15, y + 0.45, 2.6, 0.4, size=16, color=AMBER, bold=True)
    add_text(s, desc, x + 0.15, y + 0.95, 2.6, 0.9, size=12, color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: Coach Ney - Tro ly AI
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Coach Ney \u2014 Tro ly AI", "Tinh nang moi \u2014 M28")

add_text(s, "Tro ly AI ho tro giong noi, dong hanh cung hoc vien trong moi bai hoc", 0.8, 1.8, 11, 0.5, size=18, color=WHITE)

# Left - Features
add_rect(s, 0.8, 2.6, 5.8, 4.2, DARK_SURFACE)
add_text(s, "TINH NANG", 0.8, 2.7, 5.8, 0.4, size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

coach_features = [
    "Chat van ban voi AI \u2014 streaming response",
    "Ghi am giong noi qua micro (Whisper STT)",
    "AI tra loi bang giong noi (Text-to-Speech)",
    "Hieu ngu canh bai hoc/arcade hien tai",
    "Goi y nhanh (suggestion chips) de hoi AI",
    "Nhac nho khi hoc vien im lang 60s/120s (idle nudge)",
    "Nut amber phat sang goc phai man hinh",
]
add_bullet_list(s, coach_features, 1.0, 3.2, 5.4, 3.5, size=14)

# Right - Plan limits
add_rect(s, 7.2, 2.6, 5.3, 4.2, DARK_SURFACE)
add_text(s, "GIOI HAN THEO GOI", 7.2, 2.7, 5.3, 0.4, size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

plans_coach = [
    ("Basic (Mien phi)", "5 tin nhan/ngay, chi van ban"),
    ("Standard", "50 tin nhan/ngay + giong noi"),
    ("Advanced", "Khong gioi han + giong noi + VIP"),
]
y_pos = 3.3
for plan, desc in plans_coach:
    add_rect(s, 7.5, y_pos, 4.7, 0.9, DARK_CARD)
    add_text(s, plan, 7.7, y_pos + 0.05, 4.3, 0.4, size=15, color=AMBER, bold=True)
    add_text(s, desc, 7.7, y_pos + 0.45, 4.3, 0.4, size=13, color=WHITE)
    y_pos += 1.1


# ════════════════════════════════════════════════════════════════
# SLIDE 8: Cong dong hoc tap
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Cong dong hoc tap", "Tinh nang moi \u2014 M27")

add_text(s, "Thao luan va chia se kinh nghiem ngay trong bai hoc", 0.8, 1.8, 11, 0.5, size=20, color=WHITE)

disc_features = [
    "Binh luan & tra loi tren moi bai hoc Foundation va thu thach Arcade",
    "Flat threading (1 cap tra loi) \u2014 don gian, de doc",
    "Chinh sua binh luan trong 15 phut, xoa bat ky luc nao",
    "Admin kiem duyet: an/hien binh luan vi pham",
    "Rate limiting: toi da 5 binh luan/phut chong spam",
    "Bao mat XSS: lam sach moi input dau vao",
    "Avatar mau tu dong theo userId \u2014 5 bien the mau",
]
add_bullet_list(s, disc_features, 0.8, 2.6, 7, 4.0, size=16)

# Visual concept box
add_rect(s, 8.5, 2.6, 4, 4.0, DARK_SURFACE, AMBER)
add_text(s, "UX CONCEPT", 8.5, 2.7, 4, 0.4, size=14, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
concept = [
    "CommentThread",
    "\u2514\u2500 CommentCard (user, time)",
    "    \u2514\u2500 Reply cards",
    "\u2514\u2500 CommentForm",
    "",
    "Tich hop truc tiep vao",
    "trang bai hoc & arcade",
    "\u2192 Khong can roi trang",
]
add_multiline(s, concept, 8.7, 3.2, 3.6, 3.2, size=13, color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: Mo hinh kinh doanh
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Mo hinh kinh doanh")

# Three pricing cards
plans_data = [
    ("BASIC", "Mien phi", [
        "Duyet thu vien clip",
        "3 thu thach arcade/loai",
        "Video 5 giay",
        "Phan hoi AI ngan gon",
        "Coach Ney: 5 tin/ngay",
    ], False),
    ("STANDARD", "$12/thang hoac $99/nam", [
        "Tat ca noi dung Library",
        "Tat ca Arcade challenges",
        "Video 30 giay",
        "Phan hoi AI day du",
        "Coach Ney: 50 tin + giong noi",
        "Lo trinh hoc ca nhan",
    ], True),
    ("ADVANCED", "$24/thang hoac $199/nam", [
        "Tat ca tinh nang Standard",
        "Video 3 phut",
        "Bai hoc VIP doc quyen",
        "Bao cao thang tu Coach Ney",
        "Coach Ney: khong gioi han",
        "Uu tien ho tro",
    ], False),
]

for i, (name, price, features_list, highlighted) in enumerate(plans_data):
    x = 0.5 + i * 4.2
    border = AMBER if highlighted else None
    card_bg = DARK_SURFACE if not highlighted else DARK_CARD
    add_rect(s, x, 1.8, 3.9, 5.2, card_bg, border)

    if highlighted:
        add_rect(s, x, 1.8, 3.9, 0.45, AMBER)
        add_text(s, "PHO BIEN NHAT", x, 1.82, 3.9, 0.4, size=11, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, name, x, 2.35, 3.9, 0.45, size=22, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, price, x, 2.85, 3.9, 0.4, size=16, color=WHITE, align=PP_ALIGN.CENTER)

    y_feat = 3.4
    for feat in features_list:
        add_text(s, "\u2713  " + feat, x + 0.3, y_feat, 3.3, 0.35, size=13, color=WHITE)
        y_feat += 0.4

# B2B note
add_text(s, "B2B / Teams: $8/nguoi/thang  |  Enterprise: Gia tuy chinh voi SCORM/LMS", 0.5, 7.0, 12, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: Thi truong muc tieu
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Thi truong muc tieu")

# TAM/SAM/SOM
add_rect(s, 0.8, 1.8, 3.5, 2.0, DARK_SURFACE, AMBER)
add_text(s, "TAM", 0.8, 1.85, 3.5, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "$30 TY", 0.8, 2.25, 3.5, 0.5, size=32, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Thi truong dao tao chuyen nghiep\n& coaching toan cau", 0.8, 2.85, 3.5, 0.7, size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 4.8, 1.8, 3.5, 2.0, DARK_SURFACE)
add_text(s, "SAM", 4.8, 1.85, 3.5, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "$8 TY", 4.8, 2.25, 3.5, 0.5, size=32, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Nen tang ky nang chuyen nghiep\ntruyen tuyen (US + Chau Au)", 4.8, 2.85, 3.5, 0.7, size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 8.8, 1.8, 3.5, 2.0, DARK_SURFACE)
add_text(s, "SOM (3 nam)", 8.8, 1.85, 3.5, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, "$120 TRIEU", 8.8, 2.25, 3.5, 0.5, size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "100K nguoi dung tra phi\nARR trung binh $100", 8.8, 2.85, 3.5, 0.7, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Target segments
add_text(s, "DOI TUONG MUC TIEU", 0.8, 4.2, 5, 0.4, size=16, color=AMBER, bold=True)

segments = [
    ("Chuyen gia & quan ly", "Trinh bay, thuong luong, lanh dao cuoc hop"),
    ("Sinh vien & nguoi moi di lam", "Phong van xin viec, networking, giao tiep noi bo"),
    ("Dien gia & giang vien", "Dien thuyet, giang day, thuyet trinh"),
    ("Sales & CSKH", "Giao tiep voi khach hang, xu ly tinh huong"),
    ("Nguoi tim viec", "Chuan bi phong van, the hien su tu tin"),
]
y = 4.7
for title, desc in segments:
    add_text(s, "\u25CF  " + title, 0.8, y, 5, 0.35, size=15, color=WHITE, bold=True)
    add_text(s, "    " + desc, 0.8, y + 0.3, 5, 0.3, size=13, color=GREY)
    y += 0.65

# Geographic focus
add_rect(s, 7.5, 4.2, 5, 2.8, DARK_SURFACE)
add_text(s, "THI TRUONG DIA LY", 7.5, 4.3, 5, 0.4, size=16, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
geo = [
    "Giai doan 1: Viet Nam (tieng Viet)",
    "Giai doan 2: Dong Nam A (tieng Anh)",
    "Giai doan 3: Toan cau",
    "",
    "73% nha tuyen dung noi ky nang mem",
    "quan trong nhu ky nang chuyen mon",
    "(LinkedIn 2023)",
]
add_multiline(s, geo, 7.7, 4.8, 4.6, 2.0, size=14, color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 11: Lo trinh phat trien
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Lo trinh phat trien")

# Timeline rows
milestones_data = [
    ("HOAN THANH", AMBER, [
        ("M0\u2013M6", "MVP: Setup, Design, Data Pipeline, Library UI, Coaching Loop, AI Feedback, Launch"),
        ("M7\u2013M8", "Auth System + Admin CMS"),
        ("M10\u2013M13", "Script-Aware Coaching, Observation Guide, Micro-Practice, Onboarding + Learning Path"),
        ("M14\u2013M16", "100+ Lessons, Foundation Courses, Content Automation (Crawl + AI Score)"),
    ]),
    ("DA DEPLOY", GREEN, [
        ("M17\u2013M18", "Library UX Refresh + Brand Refresh (PeeTeeAI)"),
        ("M19\u2013M20", "Arcade Zone + Screenplay Crawl + ZIP Import"),
        ("M21\u2013M23", "Arcade Admin + User Analytics + Feature Analytics"),
        ("M24\u2013M25", "Access Control + Subscription Plans (PayPal + VNPay)"),
    ]),
    ("DANG PHAT TRIEN", RGBColor(0x60, 0xA5, 0xFA), [
        ("M26", "Dang ky voi phe duyet Admin"),
        ("M27", "Discussions \u2014 Binh luan & Tra loi"),
        ("M28", "Coach Ney \u2014 Tro ly AI giong noi"),
    ]),
]

y_pos = 1.8
for section_title, title_color, items in milestones_data:
    add_text(s, section_title, 0.8, y_pos, 3, 0.4, size=16, color=title_color, bold=True)
    y_pos += 0.45
    for code, desc in items:
        add_rect(s, 1.0, y_pos, 1.4, 0.38, DARK_SURFACE)
        add_text(s, code, 1.0, y_pos, 1.4, 0.38, size=11, color=title_color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, desc, 2.6, y_pos, 9.5, 0.38, size=12, color=WHITE)
        y_pos += 0.42
    y_pos += 0.2

# Note
add_text(s, "28 milestones \u2014 xay dung boi doi ngu AI-augmented trong duoi 30 ngay", 0.8, 6.8, 11, 0.4, size=14, color=GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 12: Doi ngu
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Doi ngu phat trien")

add_text(s, "Multi-Agent AI Development Team", 0.8, 1.6, 11, 0.5, size=20, color=WHITE)
add_text(s, "Mo hinh phat trien doc dao: doi ngu AI chuyen biet phoi hop qua giao thuc tin hieu", 0.8, 2.1, 11, 0.4, size=15, color=GREY)

roles = [
    ("Project Manager", "Dieu phoi milestone, quyet dinh, trang thai du an"),
    ("Designer", "UX/UI specs, design system, component specs"),
    ("Backend Engineer", "Auth, API, database, AI integration"),
    ("Data Engineer", "YouTube pipeline, screenplay parser, du lieu"),
    ("Builder", "Git, deploy Vercel, Neon DB, infrastructure"),
    ("Tester", "Test cases, bug reports, quality assurance"),
    ("Reporter", "Activity logs, documentation, knowledge base"),
    ("Marketer", "Pitch deck, brand voice, go-to-market strategy"),
]

for i, (role, desc) in enumerate(roles):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.1
    y = 2.8 + row * 2.1
    add_rect(s, x, y, 2.9, 1.7, DARK_SURFACE)
    # Role icon circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = AMBER
    circle.line.fill.background()
    add_text(s, role[0], x + 0.15, y + 0.18, 0.5, 0.45, size=20, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, role, x + 0.8, y + 0.15, 2.0, 0.4, size=13, color=AMBER, bold=True)
    add_text(s, desc, x + 0.15, y + 0.7, 2.6, 0.9, size=11, color=WHITE)

add_text(s, "Ho tro boi: PeeTeeAI JSC  |  www.peetees.ai", 0.5, 6.8, 12, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 13: Canh tranh
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)
slide_title(s, "Loi the canh tranh")

# Competition table
headers = ["Doi thu", "Ho lam gi", "Loi the cua seeneyu"]
header_widths = [2.5, 3.5, 5.5]
x_start = 0.8
y_start = 1.8

# Header row
x = x_start
for h, w in zip(headers, header_widths):
    add_rect(s, x, y_start, w, 0.55, AMBER)
    add_text(s, h, x + 0.1, y_start + 0.05, w - 0.2, 0.45, size=14, color=DARK_BG, bold=True)
    x += w

competitors = [
    ("Udemy / Coursera", "Khoa hoc video", "Xem thu dong vs Thuc hanh chu dong voi AI"),
    ("Toastmasters", "Coaching nhom truc tiep", "Co 24/7, khong can dat lich, AI-powered"),
    ("Yoodli", "AI phan hoi giong noi", "Chi giong noi; seeneyu = TOAN BO ngon ngu co the"),
    ("LinkedIn Learning", "Noi dung ky nang chuyen nghiep", "Khong vong luyen tap, khong phan hoi ca nhan"),
    ("Speeko", "App luyen noi truoc cong chung", "Chi giong noi, khong body language, khong AI Vision"),
]

for i, (comp, what, adv) in enumerate(competitors):
    y = y_start + 0.55 + i * 0.7
    bg = DARK_SURFACE if i % 2 == 0 else DARK_CARD
    x = x_start
    for val, w in zip([comp, what, adv], header_widths):
        add_rect(s, x, y, w, 0.65, bg)
        color = WHITE if val != adv else GREEN
        add_text(s, val, x + 0.1, y + 0.1, w - 0.2, 0.45, size=13, color=color)
        x += w

# Unique differentiator
add_rect(s, 0.8, 5.5, 11.5, 1.2, DARK_SURFACE, AMBER)
add_text(s, "KHAC BIET DUY NHAT", 0.8, 5.55, 11.5, 0.4, size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "seeneyu la nen tang duy nhat ket hop: (1) Trinh dien chuyen gia tu tinh huong thuc\n(2) Thuc hanh quay video chu dong  (3) Phan hoi AI da phuong thuc ve ngon ngu co the",
         1.0, 6.0, 11.0, 0.6, size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 14: Tam nhin & Lien he
# ════════════════════════════════════════════════════════════════
s = pptx.slides.add_slide(pptx.slide_layouts[6])
set_bg(s)

# Decorative amber bar at top
add_rect(s, 0, 0, 13.333, 0.08, AMBER)

add_text(s, "Tam nhin", 0.8, 0.6, 10, 0.7, size=32, color=AMBER, bold=True)
add_accent_line(s, 0.8, 1.3, 2.5)

add_text(s, "Tro thanh nen tang #1 ve huan luyen giao tiep\ntai Dong Nam A va toan cau", 0.8, 1.7, 11, 1.0, size=26, color=WHITE, bold=True)

vision_points = [
    "200+ bai hoc tu phim kinh dien & noi dung goc",
    "Ung dung mobile (iOS + Android) voi streak system",
    "B2B: Tich hop SCORM/LMS cho doanh nghiep",
    "Data flywheel: cang nhieu nguoi dung \u2192 AI cang thong minh",
    "Mo rong ra: giao tiep van hoa, phong van, dien thuyet",
]
add_bullet_list(s, vision_points, 0.8, 3.0, 8, 3.0, size=17)

# The Ask box
add_rect(s, 0.8, 5.5, 5.5, 1.5, DARK_SURFACE, AMBER)
add_text(s, "GOI VON SEED: $1.5M", 0.8, 5.55, 5.5, 0.5, size=22, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
ask_lines = [
    "40% Noi dung (mo rong 200+ clip, noi dung goc)",
    "30% San pham (mobile app, social features)",
    "20% Marketing (SEO, influencer, B2B sales)",
    "10% Van hanh (phap ly, infra, CSKH)",
]
add_multiline(s, ask_lines, 1.0, 6.1, 5.1, 1.0, size=12, color=WHITE)

# Contact box
add_rect(s, 7.0, 5.5, 5.5, 1.5, DARK_SURFACE)
add_text(s, "LIEN HE", 7.0, 5.55, 5.5, 0.5, size=18, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
contact_lines = [
    "PeeTeeAI JSC",
    "Web: www.peetees.ai",
    "App: seeneyu.vercel.app",
    "Email: contact@peetees.ai",
]
add_multiline(s, contact_lines, 7.2, 6.1, 5.1, 1.0, size=14, color=WHITE)

# Bottom tagline
add_text(s, '"Trong moi can phong quan trong \u2014 ngon ngu co the la vu khi bi mat cua ban."',
         0.8, 7.1, 12, 0.3, size=15, color=GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_dir = os.path.join(os.path.dirname(__file__), "..", "..", "..", ".shared", "outputs", "marketer")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "seeneyu-pitch-deck-vi.pptx")

pptx.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Total slides: {len(pptx.slides)}")
